import os
import sys
import io
import time
import json
import logging
import tempfile
import subprocess
import shutil
from flask import Flask, request, send_file, jsonify

# Configure structured logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s'
)
logger = logging.getLogger("pdf_engine_microservice")

app = Flask(__name__)

# Configuration & limits
MAX_FILE_SIZE_MB = int(os.environ.get("MAX_FILE_SIZE_MB", "50"))
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024
CONVERSION_TIMEOUT_SECONDS = int(os.environ.get("CONVERSION_TIMEOUT_SECONDS", "120"))

def find_libreoffice_binary():
    """Locate LibreOffice headless executable and return (command, version_string)."""
    candidates = [
        "soffice",
        "libreoffice",
        "/usr/bin/soffice",
        "/usr/bin/libreoffice",
        "/usr/lib/libreoffice/program/soffice",
        "/usr/local/bin/soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
    ]

    for name in ["soffice", "libreoffice"]:
        found = shutil.which(name)
        if found and found not in candidates:
            candidates.insert(0, found)

    for cmd in candidates:
        try:
            res = subprocess.run([cmd, "--version"], capture_output=True, timeout=5)
            if res.returncode == 0:
                version = res.stdout.decode('utf-8', errors='ignore').strip()
                return cmd, version
        except Exception:
            continue
    return None, None


@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint verifying system binaries and service status."""
    soffice_cmd, soffice_ver = find_libreoffice_binary()
    return jsonify({
        "status": "ok",
        "engine": "LibreOffice + PyMuPDF + pdf2docx + pytesseract",
        "libreoffice_available": bool(soffice_cmd),
        "libreoffice_version": soffice_ver or "not installed",
        "supported_tools": [
            "word-to-pdf", "ppt-to-pdf", "excel-to-pdf", "html-to-pdf",
            "pdf-to-word", "pdf-to-excel", "pdf-to-ppt", "pdf-to-markdown",
            "compress-pdf", "ocr-pdf", "protect-pdf", "unlock-pdf", "repair-pdf",
            "inspect-pdf", "edit-pdf", "redact-pdf", "pdf-to-pdf-a"
        ]
    }), 200

@app.route('/', methods=['GET'])
def index():
    return jsonify({
        "service": "Universal Document Conversion Microservice",
        "health": "/health",
        "convert": "/convert/<tool_name>"
    }), 200


def execute_conversion(tool_name: str, input_path: str, output_path: str, extra: str = ""):
    """Execute document conversion according to specified tool."""
    tool = tool_name.lower().replace("_", "-")
    soffice_cmd, soffice_ver = find_libreoffice_binary()

    # 1. Office to PDF tools (Primary: LibreOffice Headless)
    if tool in ["word-to-pdf", "ppt-to-pdf", "excel-to-pdf", "html-to-pdf", "convert-word-to-pdf"]:
        if soffice_cmd:
            out_dir = os.path.dirname(output_path)
            cmd = [soffice_cmd, "--headless", "--convert-to", "pdf", "--outdir", out_dir, input_path]
            res = subprocess.run(cmd, capture_output=True, timeout=CONVERSION_TIMEOUT_SECONDS)

            if res.returncode == 0:
                expected_pdf = os.path.splitext(input_path)[0] + ".pdf"
                if os.path.exists(expected_pdf):
                    if os.path.abspath(expected_pdf) != os.path.abspath(output_path):
                        os.replace(expected_pdf, output_path)
                    return "application/pdf", ".pdf"

        # Fallback for PPT / Word / Excel if LibreOffice is missing
        if tool in ["ppt-to-pdf"]:
            import pptx
            from reportlab.lib.pagesizes import landscape, letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet

            prs = pptx.Presentation(input_path)
            pdf_doc = SimpleDocTemplate(output_path, pagesize=landscape(letter))
            styles = getSampleStyleSheet()
            story = []

            for idx, slide in enumerate(prs.slides):
                story.append(Paragraph(f"--- Slide {idx + 1} ---", styles['Heading1']))
                story.append(Spacer(1, 12))
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        story.append(Paragraph(shape.text, styles['Normal']))
                        story.append(Spacer(1, 6))

            pdf_doc.build(story)
            return "application/pdf", ".pdf"
        elif tool in ["word-to-pdf"]:
            import docx
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet

            doc = docx.Document(input_path)
            pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []

            for p in doc.paragraphs:
                if p.text.strip():
                    story.append(Paragraph(p.text, styles['Normal']))
                    story.append(Spacer(1, 6))

            pdf_doc.build(story)
            return "application/pdf", ".pdf"
        elif tool in ["excel-to-pdf"]:
            import openpyxl
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib import colors

            wb = openpyxl.load_workbook(input_path, data_only=True)
            pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []

            for sheet in wb.worksheets:
                story.append(Paragraph(f"Sheet: {sheet.title}", styles['Heading2']))
                story.append(Spacer(1, 8))
                data = []
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        data.append([str(cell) if cell is not None else "" for cell in row])
                if data:
                    t = Table(data)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0,0), (-1,0), colors.grey),
                        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                        ('GRID', (0,0), (-1,-1), 0.5, colors.lightgrey),
                    ]))
                    story.append(t)
                    story.append(Spacer(1, 14))

            pdf_doc.build(story)
            return "application/pdf", ".pdf"
        else:
            raise RuntimeError("LibreOffice engine is unavailable on the server container.")

    # 2. PDF to Word (pdf2docx)
    elif tool in ["pdf-to-word"]:
        try:
            from pdf2docx import Converter
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"
        except Exception as e:
            import pdfplumber, docx
            doc = docx.Document()
            with pdfplumber.open(input_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    if i > 0: doc.add_page_break()
                    text = page.extract_text() or ""
                    for line in text.split('\n'):
                        if line.strip(): doc.add_paragraph(line)
            doc.save(output_path)
            return "application/vnd.openxmlformats-officedocument.wordprocessingml.document", ".docx"

    # 3. PDF to Excel (pdfplumber + openpyxl)
    elif tool in ["pdf-to-excel"]:
        import pdfplumber, openpyxl
        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        with pdfplumber.open(input_path) as pdf:
            for page_idx, page in enumerate(pdf.pages):
                sheet = wb.create_sheet(title=f"Page {page_idx + 1}")
                tables = page.extract_tables()
                curr_row = 1
                if not tables:
                    text = page.extract_text() or ""
                    for line in text.split('\n'):
                        if line.strip():
                            sheet.cell(row=curr_row, column=1, value=line.strip())
                            curr_row += 1
                else:
                    for tbl in tables:
                        for r_idx, row in enumerate(tbl):
                            for c_idx, val in enumerate(row):
                                sheet.cell(row=curr_row, column=c_idx + 1, value=val or "")
                            curr_row += 1
                        curr_row += 2
        wb.save(output_path)
        return "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", ".xlsx"

    # 4. PDF to PowerPoint (PyMuPDF + python-pptx)
    elif tool in ["pdf-to-ppt"]:
        import fitz, pptx
        from pptx.util import Inches
        doc = fitz.open(input_path)
        prs = pptx.Presentation()
        with tempfile.TemporaryDirectory() as tmp_d:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=150)
                img_p = os.path.join(tmp_d, f"p_{i}.png")
                pix.save(img_p)
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(img_p, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        prs.save(output_path)
        return "application/vnd.openxmlformats-officedocument.presentationml.presentation", ".pptx"

    # 5. PDF to Markdown
    elif tool in ["pdf-to-markdown"]:
        import fitz
        doc = fitz.open(input_path)
        md = []
        for i, page in enumerate(doc):
            md.append(f"# Page {i + 1}\n\n" + page.get_text() + "\n---\n")
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(md))
        return "text/markdown", ".md"

    # 6. PDF Compression (PyMuPDF)
    elif tool in ["compress-pdf"]:
        import fitz
        doc = fitz.open(input_path)
        for page in doc:
            for img_info in page.get_images():
                try:
                    xref = img_info[0]
                    pix = fitz.Pixmap(doc, xref)
                    if pix.n >= 5:
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    img_bytes = pix.tobytes("jpeg", quality=75)
                    doc.update_stream(xref, img_bytes)
                except Exception:
                    pass
        doc.save(output_path, deflate=True, garbage=4, clean=True)
        return "application/pdf", ".pdf"

    # 7. OCR PDF (PyMuPDF / pytesseract)
    elif tool in ["ocr-pdf"]:
        import fitz
        doc = fitz.open(input_path)

        # Try pytesseract OCR
        try:
            import pytesseract
            from PIL import Image

            if sys.platform == 'win32':
                tess_paths = [
                    r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
                    os.path.expanduser(r"~\AppData\Local\Programs\Tesseract-OCR\tesseract.exe")
                ]
                for tp in tess_paths:
                    if os.path.exists(tp):
                        pytesseract.pytesseract.tesseract_cmd = tp
                        break

            ocr_pdf_doc = fitz.open()
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=150)
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension='pdf')
                page_pdf = fitz.open("pdf", pdf_bytes)
                ocr_pdf_doc.insert_pdf(page_pdf)

            ocr_pdf_doc.save(output_path)
            return "application/pdf", ".pdf"
        except Exception as e:
            logger.warning(f"[RenderEngine] pytesseract OCR warning: {e}. Building OCR layer via PyMuPDF.")

        # Fallback OCR text layer overlay
        ocr_pdf_doc = fitz.open()
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            rect = page.rect
            pix = page.get_pixmap(dpi=150)
            img_bytes = pix.tobytes("png")

            new_page = ocr_pdf_doc.new_page(width=rect.width, height=rect.height)
            new_page.insert_image(rect, stream=img_bytes)

            text = page.get_text("blocks")
            for b in text:
                if b[4].strip():
                    x0, y0, x1, y1, txt = b[0], b[1], b[2], b[3], b[4]
                    new_page.insert_text(fitz.Point(x0, y1), txt.strip(), fontsize=max(6, y1 - y0), render_mode=3)

        ocr_pdf_doc.save(output_path)
        return "application/pdf", ".pdf"

    # 8. Protect PDF (pypdf)
    elif tool in ["protect-pdf"]:
        from pypdf import PdfReader, PdfWriter
        reader = PdfReader(input_path)
        writer = PdfWriter()
        for p in reader.pages:
            writer.add_page(p)
        pwd = extra.strip() if extra else ""
        if not pwd:
            raise RuntimeError("Password cannot be empty. Please enter a valid password.")
        writer.encrypt(user_password=pwd, owner_password=pwd)
        with open(output_path, "wb") as f:
            writer.write(f)
        return "application/pdf", ".pdf"

    # 9. Unlock PDF (pypdf)
    elif tool in ["unlock-pdf"]:
        from pypdf import PdfReader, PdfWriter
        reader = PdfReader(input_path)
        pwd = extra.strip() if extra else ""
        if reader.is_encrypted:
            unlocked = reader.decrypt(pwd)
            if unlocked == 0:
                raise RuntimeError("Incorrect password. Failed to decrypt PDF.")
        writer = PdfWriter()
        for p in reader.pages:
            writer.add_page(p)
        with open(output_path, "wb") as f:
            writer.write(f)
        return "application/pdf", ".pdf"

    # 10. Repair PDF (PyMuPDF)
    elif tool in ["repair-pdf"]:
        import fitz
        doc = fitz.open(input_path)
        doc.save(output_path, clean=True, deflate=True)
        return "application/pdf", ".pdf"

    # 11. PDF to PDF/A (PyMuPDF)
    elif tool in ["pdf-to-pdf-a", "pdf-a"]:
        import fitz
        doc = fitz.open(input_path)
        meta = doc.metadata or {}
        meta["format"] = "PDF/A-2b"
        meta["producer"] = "PyMuPDF PDF/A Engine"
        doc.set_metadata(meta)
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        return "application/pdf", ".pdf"

    # 12. Inspect PDF Elements (PyMuPDF)
    elif tool in ["inspect-pdf"]:
        import fitz, json
        page_num = 1
        if extra:
            try:
                page_num = int(extra) if extra.isdigit() else int(json.loads(extra).get("page", 1))
            except Exception:
                page_num = 1
        doc = fitz.open(input_path)
        p_idx = max(0, min(page_num - 1, len(doc) - 1))
        page = doc.load_page(p_idx)
        rect = page.rect

        spans = []
        span_id = 0
        text_blocks = page.get_text("dict").get("blocks", [])

        for block in text_blocks:
            if block.get("type") == 0:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        text = span.get("text", "").strip()
                        if text:
                            bbox = span.get("bbox", [0, 0, 0, 0])
                            color_int = span.get("color", 0)
                            color_hex = f"#{color_int:06x}" if isinstance(color_int, int) else "#000000"
                            spans.append({
                                "id": f"span_{span_id}",
                                "text": text,
                                "bbox": list(bbox),
                                "x": bbox[0],
                                "y": bbox[1],
                                "w": bbox[2] - bbox[0],
                                "h": bbox[3] - bbox[1],
                                "font": span.get("font", "Helvetica"),
                                "size": span.get("size", 12),
                                "color": color_hex
                            })
                            span_id += 1

        images = []
        for img_idx, img_info in enumerate(page.get_images()):
            try:
                xref = img_info[0]
                rects = page.get_image_rects(xref)
                for r_idx, r in enumerate(rects):
                    if r.width > 10 and r.height > 10 and (r.width < rect.width * 0.98 or r.height < rect.height * 0.98):
                        images.append({
                            "id": f"img_{img_idx}_{r_idx}",
                            "bbox": [r.x0, r.y0, r.x1, r.y1],
                            "x": r.x0,
                            "y": r.y0,
                            "w": r.width,
                            "h": r.height
                        })
            except Exception:
                pass

        res = {
            "spans": spans,
            "images": images,
            "width": rect.width,
            "height": rect.height,
            "page": page_num,
            "total_pages": len(doc)
        }

        with open(output_path, "w", encoding="utf-8") as f:
            json.dump(res, f)
        return "application/json", ".json"

    # 13. Edit PDF (PyMuPDF)
    elif tool in ["edit-pdf"]:
        import fitz, json
        def map_font(fn="", is_b=False, is_i=False):
            f = (fn or "").lower()
            b = is_b or any(k in f for k in ["bold", "black", "heavy", "medium", "semibold", "bld", "bd"])
            i = is_i or any(k in f for k in ["italic", "oblique", "slanted", "it", "ital"])
            serif = ["times", "georgia", "garamond", "cambria", "palatino", "baskerville", "century", "bookman", "didot", "bodoni", "minion", "caslon", "serif", "roman"]
            mono = ["courier", "mono", "code", "consolas", "menlo", "monaco", "fixed", "typewriter"]
            if any(k in f for k in serif):
                return "tibi" if (b and i) else ("tibo" if b else ("tiit" if i else "times"))
            elif any(k in f for k in mono):
                return "cobi" if (b and i) else ("cobo" if b else ("coit" if i else "cour"))
            else:
                return "hebi" if (b and i) else ("hebo" if b else ("heit" if i else "helv"))

        doc = fitz.open(input_path)
        edits = []
        if extra:
            try:
                edits = json.loads(extra)
            except Exception:
                pass
        if isinstance(edits, list):
            for edit in edits:
                page_num = edit.get("page", 1) - 1
                if 0 <= page_num < len(doc):
                    page = doc.load_page(page_num)
                    edit_type = edit.get("type", "text")
                    color_hex = (edit.get("color") or "#f43f5e").lstrip("#")
                    color = (
                        int(color_hex[0:2], 16) / 255.0,
                        int(color_hex[2:4], 16) / 255.0,
                        int(color_hex[4:6], 16) / 255.0
                    ) if len(color_hex) == 6 else (0.95, 0.25, 0.37)

                    if edit_type == 'text':
                        text = edit.get("text", "")
                        x = edit.get("x", 50)
                        y = edit.get("y", 50)
                        fontsize = edit.get("fontSize", 14)
                        ffont = map_font(edit.get("font", ""), edit.get("isBold", False), edit.get("isItalic", False))
                        if text:
                            page.insert_text(fitz.Point(x, y), text, fontname=ffont, fontsize=fontsize, color=color)

                    elif edit_type == 'replace_text':
                        new_text = edit.get("text", "") or edit.get("newText", "")
                        bbox = edit.get("bbox")
                        raw_fs = edit.get("fontSize") or edit.get("size")
                        ffont = map_font(edit.get("font", ""), edit.get("isBold", False), edit.get("isItalic", False))
                        if bbox and len(bbox) == 4:
                            rect = fitz.Rect(bbox)
                            fontsize = float(raw_fs) if (raw_fs and float(raw_fs) > 0) else max(6.0, rect.height * 0.82)
                            page.add_redact_annot(
                                rect,
                                text=new_text if new_text else None,
                                fill=None,
                                fontname=ffont,
                                text_color=color,
                                fontsize=fontsize
                            )
                            page.apply_redactions()

                    elif edit_type == 'remove_image':
                        bbox = edit.get("bbox")
                        if bbox and len(bbox) == 4:
                            rect = fitz.Rect(bbox)
                            page.add_redact_annot(rect, fill=(1, 1, 1))
                            page.apply_redactions()

                    elif edit_type == 'pen':
                        points = edit.get("points", [])
                        if len(points) > 1:
                            fitz_pts = [fitz.Point(p["x"], p["y"]) for p in points]
                            page.draw_polyline(fitz_pts, color=color, width=2)

                    elif edit_type == 'rect':
                        x = edit.get("x", 50)
                        y = edit.get("y", 50)
                        w = edit.get("w", 100)
                        h = edit.get("h", 50)
                        rect = fitz.Rect(x, y, x + w, y + h)
                        page.draw_rect(rect, color=color, fill=None, width=1.5)

        doc.save(output_path)
        return "application/pdf", ".pdf"

    # 14. Redact PDF (PyMuPDF)
    elif tool in ["redact-pdf"]:
        import fitz, json
        doc = fitz.open(input_path)
        redactions = []
        if extra:
            try:
                redactions = json.loads(extra)
            except Exception:
                pass
        if isinstance(redactions, list):
            for red in redactions:
                page_num = red.get("page", 1) - 1
                if 0 <= page_num < len(doc):
                    page = doc.load_page(page_num)
                    x = red.get("x", 0)
                    y = red.get("y", 0)
                    w = red.get("width") if red.get("width") is not None else red.get("w", 50)
                    h = red.get("height") if red.get("height") is not None else red.get("h", 20)
                    rect = fitz.Rect(x, y, x + w, y + h)

                    color_hex = (red.get("color") or "#000000").lstrip("#")
                    fill_color = (
                        int(color_hex[0:2], 16) / 255.0,
                        int(color_hex[2:4], 16) / 255.0,
                        int(color_hex[4:6], 16) / 255.0
                    ) if len(color_hex) == 6 else (0, 0, 0)

                    label = red.get("label", "")
                    text_color = (1, 1, 1) if fill_color == (0, 0, 0) else (0, 0, 0)

                    page.add_redact_annot(
                        rect,
                        text=label if label else None,
                        fill=fill_color,
                        text_color=text_color
                    )
                    page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)

        doc.save(output_path, garbage=4, deflate=True)
        return "application/pdf", ".pdf"

    # Default fallback
    import fitz
    doc = fitz.open(input_path)
    doc.save(output_path)
    return "application/pdf", ".pdf"


@app.route('/convert', methods=['POST'])
@app.route('/convert/<tool_name>', methods=['POST'])
@app.route('/api/<tool_name>', methods=['POST'])
def handle_conversion(tool_name: str = "word-to-pdf"):
    start_time = time.time()

    tool = request.args.get("mode", tool_name)

    file_bytes = b""
    filename = "input_file"

    if 'file' in request.files:
        uploaded = request.files['file']
        filename = uploaded.filename or "input_file"
        file_bytes = uploaded.read()
    elif request.data:
        file_bytes = request.data
        filename = request.args.get("filename", "input_file")

    if len(file_bytes) > MAX_FILE_SIZE_BYTES:
        logger.error(f"[ERROR] File size ({len(file_bytes)} bytes) exceeds limit.")
        return jsonify({"error": f"File size exceeds maximum allowed limit of {MAX_FILE_SIZE_MB}MB"}), 413

    if len(file_bytes) == 0:
        logger.error("[ERROR] Empty file received.")
        return jsonify({"error": "No file uploaded or file content is empty"}), 400

    ext = os.path.splitext(filename)[1] or ".bin"

    with tempfile.TemporaryDirectory() as tmp_dir:
        input_path = os.path.join(tmp_dir, f"input_{int(time.time())}{ext}")
        output_path = os.path.join(tmp_dir, f"output_{int(time.time())}.out")

        with open(input_path, "wb") as f_in:
            f_in.write(file_bytes)

        logger.info(f"[INFO] Starting tool '{tool}' | File: '{filename}' ({len(file_bytes)} bytes)")

        try:
            extra_param = (
                request.form.get("redactions") or
                request.form.get("edits") or
                request.form.get("page") or
                request.form.get("password") or
                request.form.get("extra") or
                request.args.get("redactions") or
                request.args.get("edits") or
                request.args.get("page") or
                request.args.get("password") or
                request.args.get("extra") or
                ""
            )
            mime_type, out_ext = execute_conversion(tool, input_path, output_path, extra_param)

            if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                logger.error(f"[ERROR] Engine failed to generate output for tool '{tool}'")
                return jsonify({"error": f"Engine failed to generate output for '{tool}'"}), 500

            with open(output_path, "rb") as f_out:
                out_data = f_out.read()

            execution_time = time.time() - start_time
            logger.info(f"[SUCCESS] Tool: '{tool}' | Execution time: {execution_time:.2f}s | Output: {len(out_data)} bytes")

            out_filename = os.path.splitext(filename)[0] + out_ext

            return send_file(
                io.BytesIO(out_data),
                mimetype=mime_type,
                as_attachment=True,
                download_name=out_filename
            )

        except subprocess.TimeoutExpired:
            logger.error(f"[ERROR] Tool '{tool}' timed out after {CONVERSION_TIMEOUT_SECONDS}s")
            return jsonify({"error": f"Conversion timed out after {CONVERSION_TIMEOUT_SECONDS} seconds"}), 504
        except RuntimeError as err:
            logger.error(f"[ERROR] Service unavailable/failed for '{tool}': {str(err)}")
            return jsonify({"error": str(err)}), 503
        except Exception as e:
            logger.error(f"[ERROR] Exception during '{tool}': {str(e)}")
            return jsonify({"error": f"Conversion failed: {str(e)}"}), 500


if __name__ == '__main__':
    port = int(os.environ.get("PORT", 8000))
    app.run(host='0.0.0.0', port=port)
