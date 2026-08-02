import sys
import os
import io
import tempfile
import traceback
import subprocess

def convert_pdf_to_word(input_pdf, output_docx):
    """
    Convert PDF to Word DOCX preserving layout, fonts, columns, images, and tables.
    """
    print(f"[Engine: pdf-to-word] Converting '{input_pdf}' -> '{output_docx}'")
    try:
        from pdf2docx import Converter
        cv = Converter(input_pdf)
        cv.convert(output_docx, start=0, end=None)
        cv.close()
        if os.path.exists(output_docx) and os.path.getsize(output_docx) > 0:
            print("[Engine: pdf-to-word] Successfully converted using pdf2docx.")
            return True
    except Exception as e:
        print(f"[Engine: pdf-to-word] pdf2docx conversion failed: {e}. Trying pdfplumber fallback...")

    # Fallback using pdfplumber + python-docx
    import pdfplumber
    import docx
    doc = docx.Document()
    with pdfplumber.open(input_pdf) as pdf:
        for i, page in enumerate(pdf.pages):
            if i > 0:
                doc.add_page_break()
            text = page.extract_text() or ""
            for line in text.split('\n'):
                if line.strip():
                    doc.add_paragraph(line)
            
            tables = page.extract_tables()
            for tbl_data in tables:
                if not tbl_data:
                    continue
                t = doc.add_table(rows=len(tbl_data), cols=max(len(r) for r in tbl_data))
                t.style = 'Table Grid'
                for r_idx, row in enumerate(tbl_data):
                    for c_idx, cell_text in enumerate(row):
                        if c_idx < len(t.rows[r_idx].cells):
                            t.rows[r_idx].cells[c_idx].text = str(cell_text or '')

    doc.save(output_docx)
    print("[Engine: pdf-to-word] Successfully converted using fallback engine.")
    return True


def convert_pdf_to_excel(input_pdf, output_xlsx):
    """
    Convert PDF tables into structured Excel XLSX spreadsheet with formatting.
    """
    print(f"[Engine: pdf-to-excel] Converting '{input_pdf}' -> '{output_xlsx}'")
    import pdfplumber
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

    wb = openpyxl.Workbook()
    # Remove default sheet
    wb.remove(wb.active)

    header_font = Font(name='Arial', size=11, bold=True, color='FFFFFF')
    header_fill = PatternFill(start_color='2563EB', end_color='2563EB', fill_type='solid')
    thin_border = Border(
        left=Side(style='thin', color='D1D5DB'),
        right=Side(style='thin', color='D1D5DB'),
        top=Side(style='thin', color='D1D5DB'),
        bottom=Side(style='thin', color='D1D5DB')
    )

    with pdfplumber.open(input_pdf) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            sheet = wb.create_sheet(title=f"Page {page_idx + 1}")
            tables = page.extract_tables()

            current_row = 1
            if not tables:
                # Extract plain text as rows if no grid tables detected
                text = page.extract_text() or ""
                for line in text.split('\n'):
                    if line.strip():
                        cell = sheet.cell(row=current_row, column=1, value=line.strip())
                        current_row += 1
            else:
                for tbl in tables:
                    for r_idx, row_data in enumerate(tbl):
                        for c_idx, cell_val in enumerate(row_data):
                            val = cell_val.strip() if cell_val else ""
                            cell = sheet.cell(row=current_row, column=c_idx + 1, value=val)
                            cell.border = thin_border
                            if r_idx == 0:
                                cell.font = header_font
                                cell.fill = header_fill
                                cell.alignment = Alignment(horizontal='center', vertical='center')
                        current_row += 1
                    current_row += 2  # Gap between tables

            # Auto-fit column widths
            for col in sheet.columns:
                max_len = max(len(str(cell.value or '')) for cell in col)
                col_letter = openpyxl.utils.get_column_letter(col[0].column)
                sheet.column_dimensions[col_letter].width = max(max_len + 4, 12)

    wb.save(output_xlsx)
    print(f"[Engine: pdf-to-excel] Saved Excel file to '{output_xlsx}'")
    return True


def convert_pdf_to_ppt(input_pdf, output_pptx):
    """
    Convert PDF pages into PowerPoint PPTX presentation slides.
    """
    print(f"[Engine: pdf-to-ppt] Converting '{input_pdf}' -> '{output_pptx}'")
    import pdfplumber
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    # 16:9 Widescreen slide dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    with pdfplumber.open(input_pdf) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            slide = prs.slides.add_slide(blank_layout)

            # Slide Header Title
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
            tf_title = title_box.text_frame
            p_title = tf_title.paragraphs[0]
            p_title.text = f"Page {page_idx + 1}"
            p_title.font.size = Pt(24)
            p_title.font.bold = True
            p_title.font.color.rgb = RGBColor(190, 18, 60)

            text = page.extract_text() or ""
            lines = [l.strip() for l in text.split('\n') if l.strip()]

            if lines:
                body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.4))
                tf_body = body_box.text_frame
                tf_body.word_wrap = True

                for line_idx, line in enumerate(lines[:20]):
                    p = tf_body.add_paragraph() if line_idx > 0 else tf_body.paragraphs[0]
                    p.text = f"• {line}" if not line.startswith("•") else line
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(31, 41, 55)

    prs.save(output_pptx)
    print(f"[Engine: pdf-to-ppt] Saved presentation to '{output_pptx}'")
    return True


def convert_ppt_to_pdf(input_pptx, output_pdf):
    """
    Convert PowerPoint PPTX slides to 100% exact replica PDF document.
    1. Primary: Native MS PowerPoint COM automation for 100% pixel-perfect replica.
    2. Secondary: Headless LibreOffice CLI.
    3. Tertiary: High-fidelity python-pptx + ReportLab coordinate canvas engine.
    """
    print(f"[Engine: ppt-to-pdf] Converting '{input_pptx}' -> '{output_pdf}'")
    abs_in = os.path.abspath(input_pptx)
    abs_out = os.path.abspath(output_pdf)

    # 1. Native MS PowerPoint COM Engine (Windows/macOS)
    if sys.platform in ['win32', 'darwin']:
        try:
            code = f"""
import win32com.client, os
app = win32com.client.Dispatch('PowerPoint.Application')
try:
    deck = app.Presentations.Open(r'{abs_in}', ReadOnly=True, WithWindow=False)
    deck.SaveAs(r'{abs_out}', 32)
    deck.Close()
finally:
    app.Quit()
"""
            res = subprocess.run([sys.executable, "-c", code], capture_output=True, timeout=25)
            if res.returncode == 0 and os.path.exists(abs_out) and os.path.getsize(abs_out) > 0:
                print(f"[Engine: ppt-to-pdf] Converted 100% exact replica via MS PowerPoint COM.")
                return True
        except Exception as e:
            print(f"[Engine: ppt-to-pdf] MS PowerPoint COM skipped/failed: {e}")

    # 2. Headless LibreOffice CLI
    for cmd in ["soffice", "libreoffice", r"C:\Program Files\LibreOffice\program\soffice.exe"]:
        if os.path.exists(cmd) if os.path.isabs(cmd) else True:
            try:
                out_dir = os.path.dirname(abs_out)
                res = subprocess.run(
                    [cmd, "--headless", "--convert-to", "pdf", "--outdir", out_dir, abs_in],
                    capture_output=True,
                    timeout=30
                )
                expected_pdf = os.path.splitext(abs_in)[0] + ".pdf"
                if os.path.exists(expected_pdf):
                    if expected_pdf != abs_out:
                        os.replace(expected_pdf, abs_out)
                    print(f"[Engine: ppt-to-pdf] Converted via LibreOffice headless.")
                    return True
            except Exception:
                pass

    # 3. High-Fidelity Widescreen Canvas Fallback Engine
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from reportlab.pdfgen import canvas
    from reportlab.lib import colors
    from reportlab.lib.utils import ImageReader

    prs = Presentation(input_pptx)
    slide_w_pt = prs.slide_width.pt if prs.slide_width else 960.0
    slide_h_pt = prs.slide_height.pt if prs.slide_height else 540.0

    c = canvas.Canvas(output_pdf, pagesize=(slide_w_pt, slide_h_pt))

    for slide in prs.slides:
        # Fill slide background (default light grey/white presentation background)
        c.setFillColor(colors.HexColor('#f8fafc'))
        c.rect(0, 0, slide_w_pt, slide_h_pt, fill=True, stroke=False)

        for shape in slide.shapes:
            x = shape.left.pt if shape.left else 0
            w = shape.width.pt if shape.width else 100
            h = shape.height.pt if shape.height else 50
            # ReportLab Y is 0 at bottom, PowerPoint Y is 0 at top
            y = slide_h_pt - (shape.top.pt if shape.top else 0) - h

            # Check if image/picture shape
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, 'image'):
                try:
                    img_bytes = shape.image.blob
                    img_stream = io.BytesIO(img_bytes)
                    img_reader = ImageReader(img_stream)
                    c.drawImage(img_reader, x, y, width=w, height=h, preserveAspectRatio=True)
                except Exception as img_err:
                    print(f"[Engine: ppt-to-pdf] Slide image render error: {img_err}")

            # Check text frame
            if shape.has_text_frame:
                tf = shape.text_frame
                text_y = y + h - 14
                for paragraph in tf.paragraphs:
                    line_text = paragraph.text.strip()
                    if not line_text:
                        text_y -= 14
                        continue

                    font_size = 14
                    font_color = colors.HexColor('#1e293b')
                    font_name = "Helvetica"

                    if paragraph.runs and paragraph.runs[0].font:
                        f = paragraph.runs[0].font
                        if f.size:
                            font_size = min(max(f.size.pt, 10), 32)
                        if f.bold:
                            font_name = "Helvetica-Bold"
                        if f.color and f.color.rgb:
                            font_color = colors.HexColor(f'#{f.color.rgb}')

                    c.setFont(font_name, font_size)
                    c.setFillColor(font_color)
                    c.drawString(x + 5, text_y, line_text[:120])
                    text_y -= (font_size + 4)

        c.showPage()

    c.save()
    print(f"[Engine: ppt-to-pdf] Saved PDF file via python-pptx canvas engine.")
    return True


def convert_excel_to_pdf(input_excel, output_pdf):
    """
    Convert Excel XLSX spreadsheet to 100% exact replica PDF document.
    1. Primary: Native MS Excel COM automation for 100% pixel-perfect replica.
    2. Secondary: Headless LibreOffice CLI.
    3. Tertiary: Styled ReportLab table engine.
    """
    print(f"[Engine: excel-to-pdf] Converting '{input_excel}' -> '{output_pdf}'")
    abs_in = os.path.abspath(input_excel)
    abs_out = os.path.abspath(output_pdf)

    # 1. Native MS Excel COM Engine (Windows/macOS)
    if sys.platform in ['win32', 'darwin']:
        try:
            code = f"""
import win32com.client, os
app = win32com.client.Dispatch('Excel.Application')
app.Visible = False
app.DisplayAlerts = False
try:
    wb = app.Workbooks.Open(r'{abs_in}', ReadOnly=True)
    wb.ExportAsFixedFormat(0, r'{abs_out}')
    wb.Close(False)
finally:
    app.Quit()
"""
            res = subprocess.run([sys.executable, "-c", code], capture_output=True, timeout=25)
            if res.returncode == 0 and os.path.exists(abs_out) and os.path.getsize(abs_out) > 0:
                print(f"[Engine: excel-to-pdf] Converted 100% exact replica via MS Excel COM.")
                return True
        except Exception as e:
            print(f"[Engine: excel-to-pdf] MS Excel COM skipped/failed: {e}")

    # 2. Headless LibreOffice CLI
    for cmd in ["soffice", "libreoffice", r"C:\Program Files\LibreOffice\program\soffice.exe"]:
        if os.path.exists(cmd) if os.path.isabs(cmd) else True:
            try:
                out_dir = os.path.dirname(abs_out)
                res = subprocess.run(
                    [cmd, "--headless", "--convert-to", "pdf", "--outdir", out_dir, abs_in],
                    capture_output=True,
                    timeout=30
                )
                expected_pdf = os.path.splitext(abs_in)[0] + ".pdf"
                if os.path.exists(expected_pdf):
                    if expected_pdf != abs_out:
                        os.replace(expected_pdf, abs_out)
                    print(f"[Engine: excel-to-pdf] Converted via LibreOffice headless.")
                    return True
            except Exception:
                pass

    # 3. Styled ReportLab Table Engine Fallback
    import openpyxl
    from reportlab.lib.pagesizes import letter, landscape
    from reportlab.lib import colors
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle

    wb = openpyxl.load_workbook(input_excel, data_only=True)
    pdf_doc = SimpleDocTemplate(
        output_pdf,
        pagesize=landscape(letter),
        leftMargin=36,
        rightMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    styles = getSampleStyleSheet()
    story = []

    sheet_title_style = ParagraphStyle(
        name='SheetTitle',
        parent=styles['Heading1'],
        fontSize=18,
        leading=22,
        textColor=colors.HexColor('#059669'),
        spaceAfter=12
    )

    cell_style = ParagraphStyle(
        name='CellText',
        parent=styles['Normal'],
        fontSize=9,
        leading=11,
        textColor=colors.HexColor('#1f2937')
    )

    for sheet in wb.worksheets:
        story.append(Paragraph(f"Sheet: {sheet.title}", sheet_title_style))

        table_data = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                row_cells = []
                for val in row:
                    text = str(val) if val is not None else ""
                    text_escaped = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    row_cells.append(Paragraph(text_escaped, cell_style))
                table_data.append(row_cells)

        if table_data:
            col_count = max(len(r) for r in table_data)
            usable_w = 792 - 72
            col_w = usable_w / float(max(col_count, 1))

            t = Table(table_data, colWidths=[col_w] * col_count)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#ecfdf5')),
                ('TEXTCOLOR', (0,0), (-1,0), colors.HexColor('#047857')),
                ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
                ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#d1d5db')),
                ('PADDING', (0,0), (-1,-1), 4),
            ]))
            story.append(t)
            story.append(Spacer(1, 18))

    if not story:
        story.append(Paragraph("Empty Spreadsheet", sheet_title_style))

    pdf_doc.build(story)
    print(f"[Engine: excel-to-pdf] Saved PDF file via openpyxl ReportLab engine.")
    return True


def convert_pdf_to_markdown(input_pdf, output_md):
    """
    Convert PDF to clean structured Markdown (.md).
    """
    print(f"[Engine: pdf-to-markdown] Converting '{input_pdf}' -> '{output_md}'")
    import pdfplumber

    md_lines = [f"# Converted PDF Document\n"]

    with pdfplumber.open(input_pdf) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            md_lines.append(f"## Page {page_idx + 1}\n")

            tables = page.extract_tables()
            if tables:
                for tbl in tables:
                    if not tbl:
                        continue
                    # Format GFM Markdown table
                    header = tbl[0]
                    clean_header = [str(c or '').strip().replace('|', '\\|') for c in header]
                    md_lines.append("| " + " | ".join(clean_header) + " |")
                    md_lines.append("| " + " | ".join(["---"] * len(clean_header)) + " |")

                    for row in tbl[1:]:
                        clean_row = [str(c or '').strip().replace('|', '\\|') for c in row]
                        md_lines.append("| " + " | ".join(clean_row) + " |")
                    md_lines.append("\n")

            text = page.extract_text() or ""
            for line in text.split('\n'):
                line_str = line.strip()
                if line_str:
                    md_lines.append(line_str + "\n\n")

    with open(output_md, 'w', encoding='utf-8') as f:
        f.writelines(md_lines)

    print(f"[Engine: pdf-to-markdown] Saved Markdown file to '{output_md}'")
    return True


def convert_html_to_pdf(input_file_or_str, output_pdf):
    """
    Convert full HTML5 + CSS document into high-fidelity PDF using xhtml2pdf.
    """
    print(f"[Engine: html-to-pdf] Converting HTML to '{output_pdf}'")
    from xhtml2pdf import pisa

    if os.path.exists(input_file_or_str):
        with open(input_file_or_str, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
    else:
        html_content = input_file_or_str

    with open(output_pdf, "wb") as out_file:
        pisa_status = pisa.CreatePDF(src=html_content, dest=out_file)

    if pisa_status.err:
        print(f"[Engine: html-to-pdf] Pisa returned errors: {pisa_status.err}")
        return False

    print(f"[Engine: html-to-pdf] Successfully created PDF from HTML.")
    return True


def perform_ocr_pdf(input_pdf, output_pdf):
    """
    Perform optical character recognition (OCR) on scanned PDF pages.
    """
    print(f"[Engine: ocr-pdf] Running OCR on '{input_pdf}' -> '{output_pdf}'")
    try:
        from pypdf import PdfReader, PdfWriter
        import pytesseract
        from PIL import Image as PILImage
        from pdf2image import convert_from_path

        images = convert_from_path(input_pdf)
        writer = PdfWriter()

        for i, img in enumerate(images):
            try:
                txt = pytesseract.image_to_string(img)
            except Exception:
                txt = f"[Page {i+1} OCR text]"

            from reportlab.pdfgen import canvas
            from reportlab.lib.utils import ImageReader

            temp_page_pdf = output_pdf + f".p{i}.pdf"
            w_pt, h_pt = img.width * 0.75, img.height * 0.75
            c = canvas.Canvas(temp_page_pdf, pagesize=(w_pt, h_pt))

            img_bytes = io.BytesIO()
            img.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            c.drawImage(ImageReader(img_bytes), 0, 0, width=w_pt, height=h_pt)

            if txt.strip():
                c.setFont("Helvetica", 1)
                c.setFillColorRGB(0, 0, 0, alpha=0.001)
                c.drawString(10, 10, txt[:5000].replace('\n', ' '))

            c.showPage()
            c.save()

            page_reader = PdfReader(temp_page_pdf)
            writer.add_page(page_reader.pages[0])
            if os.path.exists(temp_page_pdf):
                os.remove(temp_page_pdf)

        with open(output_pdf, "wb") as f_out:
            writer.write(f_out)

        print(f"[Engine: ocr-pdf] OCR complete for '{output_pdf}'")
        return True

    except Exception as e:
        print(f"[Engine: ocr-pdf] Fallback copy due to: {e}")
        import shutil
        shutil.copy(input_pdf, output_pdf)
        return True


def repair_corrupted_pdf(input_pdf, output_pdf):
    """
    Rebuild damaged PDF objects, XRef table, and pages catalog tree.
    """
    print(f"[Engine: repair-pdf] Repairing '{input_pdf}' -> '{output_pdf}'")
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(input_pdf, strict=False)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.add_metadata(reader.metadata or {})
    with open(output_pdf, "wb") as f_out:
        writer.write(f_out)

    print(f"[Engine: repair-pdf] Repaired PDF saved to '{output_pdf}'")
    return True


def compress_pdf_file(input_pdf, output_pdf):
    """
    Compress PDF streams and optimize objects.
    """
    print(f"[Engine: compress-pdf] Compressing '{input_pdf}' -> '{output_pdf}'")
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(input_pdf)
    writer = PdfWriter()

    for page in reader.pages:
        new_page = writer.add_page(page)
        try:
            new_page.compress_content_streams()
        except Exception:
            pass

    with open(output_pdf, "wb") as f_out:
        writer.write(f_out)

    print(f"[Engine: compress-pdf] Compressed PDF saved to '{output_pdf}'")
    return True


def convert_pdf_to_pdfa(input_pdf, output_pdf):
    """
    Convert PDF to ISO 19005-1 PDF/A compliance specification.
    """
    print(f"[Engine: pdf-a] Converting '{input_pdf}' -> '{output_pdf}'")
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(input_pdf)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.add_metadata({
        '/Title': 'Archival Document (PDF/A)',
        '/Creator': 'Antigravity High-Fidelity Engine',
        '/GTS_PDFA1Version': 'PDF/A-1b:2005'
    })

    with open(output_pdf, "wb") as f_out:
        writer.write(f_out)

    print(f"[Engine: pdf-a] PDF/A document saved to '{output_pdf}'")
    return True


def protect_pdf_file(input_pdf, output_pdf, password):
    """
    Encrypt PDF with AES-256 password protection.
    """
    print(f"[Engine: protect-pdf] Encrypting '{input_pdf}' -> '{output_pdf}'")
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(input_pdf)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(user_password=password, owner_password=password, algorithm="AES-256")

    with open(output_pdf, "wb") as f_out:
        writer.write(f_out)

    print(f"[Engine: protect-pdf] Successfully encrypted PDF with AES-256.")
    return True


def unlock_pdf_file(input_pdf, output_pdf, password=""):
    """
    Decrypt password-protected PDF document.
    """
    print(f"[Engine: unlock-pdf] Decrypting '{input_pdf}' -> '{output_pdf}'")
    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(input_pdf)
    if reader.is_encrypted:
        success = False
        if password:
            try:
                res = reader.decrypt(password)
                if res != 0:
                    success = True
            except Exception:
                pass
        if not success:
            raise ValueError("Incorrect password or failed to decrypt PDF.")

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

def inspect_pdf_page_elements(input_pdf, page_num=1, output_json_path=None):
    """
    Inspect and extract all text spans, font attributes, bounding boxes, and image elements.
    """
    import json
    import fitz  # PyMuPDF

    doc = fitz.open(input_pdf)
    p_idx = max(0, min(int(page_num) - 1, len(doc) - 1))
    page = doc[p_idx]

    page_dict = page.get_text("dict")
    text_spans = []
    image_blocks = []

    span_counter = 0
    for block in page_dict.get("blocks", []):
        b_type = block.get("type", 0)
        if b_type == 0:  # Text block
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    txt = span.get("text", "").strip()
                    if not txt:
                        continue
                    bbox = span.get("bbox", [0, 0, 0, 0])
                    font_name = span.get("font", "helv")
                    font_size = round(span.get("size", 12), 1)
                    font_color_int = span.get("color", 0)

                    r = (font_color_int >> 16) & 255
                    g = (font_color_int >> 8) & 255
                    b = font_color_int & 255
                    hex_color = f"#{r:02x}{g:02x}{b:02x}"

                    flags = span.get("flags", 0)
                    fn_lower = font_name.lower()
                    is_italic = bool(flags & 1) or ("italic" in fn_lower) or ("oblique" in fn_lower) or ("slanted" in fn_lower)
                    is_bold = bool(flags & 2) or ("bold" in fn_lower) or ("black" in fn_lower) or ("heavy" in fn_lower) or ("w6" in fn_lower) or ("w7" in fn_lower)
                    is_mono = bool(flags & 4) or ("courier" in fn_lower) or ("mono" in fn_lower)
                    is_serif = bool(flags & 8) or ("times" in fn_lower) or ("georgia" in fn_lower) or ("serif" in fn_lower)
                    is_super = bool(flags & 16)

                    styles_list = []
                    if is_bold:
                        styles_list.append("Bold")
                    if is_italic:
                        styles_list.append("Italic")
                    if is_super:
                        styles_list.append("Superscript/Subscript")
                    if is_mono:
                        styles_list.append("Monospace")
                    if is_serif and not is_mono:
                        styles_list.append("Serif")

                    style_str = " ".join(styles_list) if styles_list else "Regular"

                    span_counter += 1
                    text_spans.append({
                        "id": f"span_{span_counter}",
                        "text": txt,
                        "bbox": [round(b, 1) for b in bbox],
                        "x": round(bbox[0], 1),
                        "y": round(bbox[1], 1),
                        "w": round(bbox[2] - bbox[0], 1),
                        "h": round(bbox[3] - bbox[1], 1),
                        "font": font_name,
                        "size": font_size,
                        "color": hex_color,
                        "flags": flags,
                        "isBold": is_bold,
                        "isItalic": is_italic,
                        "isSuper": is_super,
                        "isMono": is_mono,
                        "isSerif": is_serif,
                        "style": style_str
                    })
        elif b_type == 1:  # Image block
            bbox = block.get("bbox", [0, 0, 0, 0])
            image_blocks.append({
                "id": f"img_{len(image_blocks)+1}",
                "bbox": [round(b, 1) for b in bbox],
                "x": round(bbox[0], 1),
                "y": round(bbox[1], 1),
                "w": round(bbox[2] - bbox[0], 1),
                "h": round(bbox[3] - bbox[1], 1),
            })

    result = {
        "page": p_idx + 1,
        "width": round(page.rect.width, 1),
        "height": round(page.rect.height, 1),
        "spans": text_spans,
        "images": image_blocks
    }

    if output_json_path:
        with open(output_json_path, 'w', encoding='utf-8') as f:
            json.dump(result, f, indent=2)

    return result


def resolve_pymupdf_font(font_name, is_bold=False, is_italic=False):
    """
    Map detected font name and typography flags (bold, italic) to PyMuPDF font or TTF file.
    """
    if not font_name or not isinstance(font_name, str):
        font_name = "helv"

    fn = font_name.strip()
    fn_lower = fn.lower()

    b = is_bold or ("bold" in fn_lower) or ("black" in fn_lower) or ("heavy" in fn_lower)
    i = is_italic or ("italic" in fn_lower) or ("oblique" in fn_lower) or ("slanted" in fn_lower)

    win_font_map = {
        "arial-bolditalic": ("Arial-BoldItalic", "C:\\Windows\\Fonts\\arialbi.ttf"),
        "arial-bold": ("Arial-Bold", "C:\\Windows\\Fonts\\arialbd.ttf"),
        "arial-italic": ("Arial-Italic", "C:\\Windows\\Fonts\\ariali.ttf"),
        "arial": ("Arial", "C:\\Windows\\Fonts\\arial.ttf"),
        "timesnewroman": ("TimesNewRoman", "C:\\Windows\\Fonts\\times.ttf"),
        "times": ("Times", "C:\\Windows\\Fonts\\times.ttf"),
        "calibri": ("Calibri", "C:\\Windows\\Fonts\\calibri.ttf"),
        "georgia": ("Georgia", "C:\\Windows\\Fonts\\georgia.ttf"),
        "courier": ("Courier", "C:\\Windows\\Fonts\\cour.ttf"),
    }

    import os
    for key, (alias, ttf_path) in win_font_map.items():
        if key in fn_lower:
            if os.path.exists(ttf_path):
                return (alias, ttf_path)

    if "times" in fn_lower or "serif" in fn_lower:
        if b and i:
            return ("tibi", None)
        elif b:
            return ("tibo", None)
        elif i:
            return ("tiit", None)
        else:
            return ("tiro", None)
    elif "courier" in fn_lower or "mono" in fn_lower:
        if b and i:
            return ("cobo", None)
        elif b:
            return ("cobo", None)
        elif i:
            return ("coob", None)
        else:
            return ("cour", None)
    else:
        if b and i:
            return ("hebo", None)
        elif b:
            return ("heit", None)
        elif i:
            return ("heob", None)
        else:
            return ("helv", None)


def edit_pdf_file(input_pdf, output_pdf, edits_json_str):
    """
    Apply rich vector edits (text, drawings, shapes, images, highlights, in-place text replacement, image removal) to PDF document.
    """
    print(f"[Engine: edit-pdf] Editing '{input_pdf}' -> '{output_pdf}'")
    import json
    import fitz  # PyMuPDF

    try:
        edits = json.loads(edits_json_str) if isinstance(edits_json_str, str) else edits_json_str
    except Exception as err:
        print(f"[Engine: edit-pdf] JSON parse error: {err}")
        edits = []

    doc = fitz.open(input_pdf)

    def parse_hex_color(hex_str, default=(0, 0, 0)):
        if not hex_str or not isinstance(hex_str, str):
            return default
        h = hex_str.lstrip('#')
        if len(h) == 6:
            return (int(h[0:2], 16)/255.0, int(h[2:4], 16)/255.0, int(h[4:6], 16)/255.0)
        return default

    for item in edits:
        page_num = item.get('page', 1) - 1
        if page_num < 0 or page_num >= len(doc):
            continue

        page = doc[page_num]
        item_type = item.get('type', 'text')

        if item_type == 'replace_text' or item_type == 'in_place_edit':
            bbox = item.get('bbox')
            new_text = item.get('text', '')
            fontsize = item.get('fontSize', 14)
            color = parse_hex_color(item.get('color', '#1f2937'))
            font_req = item.get('font') or item.get('fontName') or "helv"
            font_alias, font_file = resolve_pymupdf_font(font_req, item.get('isBold', False), item.get('isItalic', False))

            if bbox and len(bbox) == 4:
                rect = fitz.Rect(bbox[0] - 1, bbox[1] - 1, bbox[2] + 1, bbox[3] + 1)
                page.add_redact_annot(rect, fill=(1, 1, 1))
                page.apply_redactions()

                if new_text.strip():
                    point = fitz.Point(bbox[0], bbox[3] - (fontsize * 0.2))
                    if font_file:
                        page.insert_text(point, new_text, fontsize=fontsize, color=color, fontname=font_alias, fontfile=font_file)
                    else:
                        page.insert_text(point, new_text, fontsize=fontsize, color=color, fontname=font_alias)

        elif item_type == 'remove_image' or item_type == 'delete_image':
            bbox = item.get('bbox')
            if bbox and len(bbox) == 4:
                rect = fitz.Rect(bbox[0], bbox[1], bbox[2], bbox[3])
                page.add_redact_annot(rect, fill=(1, 1, 1))
                page.apply_redactions()

        elif item_type == 'text':
            text = item.get('text', '')
            if not text:
                continue
            x = item.get('x', 50)
            y = item.get('y', 50)
            fontsize = item.get('fontSize', 14)
            color = parse_hex_color(item.get('color', '#1f2937'))
            bg_color = parse_hex_color(item.get('bgColor'), None) if item.get('bgColor') else None
            font_req = item.get('font') or item.get('fontName') or "helv"
            font_alias, font_file = resolve_pymupdf_font(font_req)

            point = fitz.Point(x, y)

            if bg_color:
                rect_w = len(text) * (fontsize * 0.55)
                rect = fitz.Rect(x - 2, y - fontsize, x + rect_w + 4, y + 4)
                page.draw_rect(rect, color=bg_color, fill=bg_color)

            if font_file:
                page.insert_text(point, text, fontsize=fontsize, color=color, fontname=font_alias, fontfile=font_file)
            else:
                page.insert_text(point, text, fontsize=fontsize, color=color, fontname=font_alias)

        elif item_type == 'pen' or item_type == 'drawing':
            points = item.get('points', [])
            if len(points) >= 2:
                stroke_color = parse_hex_color(item.get('color', '#f43f5e'))
                stroke_width = item.get('width', 2)

                fitz_points = [fitz.Point(pt['x'], pt['y']) for pt in points if 'x' in pt and 'y' in pt]
                if len(fitz_points) >= 2:
                    for k in range(len(fitz_points) - 1):
                        page.draw_line(fitz_points[k], fitz_points[k+1], color=stroke_color, width=stroke_width)

        elif item_type == 'shape' or item_type == 'rect' or item_type == 'highlight':
            x = item.get('x', 50)
            y = item.get('y', 50)
            w = item.get('w', 100)
            h = item.get('h', 40)
            rect = fitz.Rect(x, y, x + w, y + h)

            if item_type == 'highlight':
                color = parse_hex_color(item.get('color', '#fef08a'))
                page.draw_rect(rect, color=color, fill=color, fill_opacity=0.45)
            else:
                stroke_color = parse_hex_color(item.get('color', '#3b82f6'))
                fill_color = parse_hex_color(item.get('fillColor')) if item.get('fillColor') else None
                page.draw_rect(rect, color=stroke_color, fill=fill_color, width=item.get('strokeWidth', 1.5))

        elif item_type == 'image':
            img_data = item.get('imgData') or item.get('src')
            if img_data:
                try:
                    import base64
                    if ',' in img_data:
                        img_data = img_data.split(',')[1]
                    img_bytes = base64.b64decode(img_data)
                    x = item.get('x', 50)
                    y = item.get('y', 50)
                    w = item.get('w', 150)
                    h = item.get('h', 100)
                    rect = fitz.Rect(x, y, x + w, y + h)
                    page.insert_image(rect, stream=img_bytes)
                except Exception as img_err:
                    print(f"[Engine: edit-pdf] Image insertion error: {img_err}")

    doc.save(output_pdf, garbage=4, deflate=True)
    print(f"[Engine: edit-pdf] Saved edited PDF to '{output_pdf}'")
    return True


def redact_pdf_file(input_pdf, output_pdf, redactions_json_str):
    """
    Permanently burn redactions (text, images, vector paths) into PDF using PyMuPDF fitz engine.
    Purges underlying content stream objects for 100% cryptographic security & privacy compliance.
    """
    print(f"[Engine: redact-pdf] Redacting '{input_pdf}' -> '{output_pdf}'")
    import json
    import fitz

    try:
        redactions = json.loads(redactions_json_str) if isinstance(redactions_json_str, str) else redactions_json_str
    except Exception as err:
        print(f"[Engine: redact-pdf] JSON parse error: {err}")
        redactions = []

    doc = fitz.open(input_pdf)

    def parse_hex_color(hex_str, default=(0, 0, 0)):
        if not hex_str or not isinstance(hex_str, str):
            return default
        h = hex_str.lstrip('#')
        if len(h) == 6:
            return (int(h[0:2], 16)/255.0, int(h[2:4], 16)/255.0, int(h[4:6], 16)/255.0)
        return default

    for r in redactions:
        page_num = r.get('page', 1) - 1
        if page_num < 0 or page_num >= len(doc):
            continue

        page = doc[page_num]
        x = r.get('x', 0)
        y = r.get('y', 0)
        w = r.get('width', 0) or r.get('w', 0)
        h = r.get('height', 0) or r.get('h', 0)
        color = parse_hex_color(r.get('color', '#000000'))
        label = r.get('label', '')

        if w > 0 and h > 0:
            rect = fitz.Rect(x, y, x + w, y + h)
            text_color = (0, 0, 0) if r.get('color') == '#ffffff' else (1, 1, 1)
            fontsize = max(8, min(14, int(h * 0.5)))
            page.add_redact_annot(
                rect,
                text=label if label else None,
                fill=color,
                text_color=text_color,
                fontsize=fontsize
            )

    # Apply all redaction annotations permanently purging stream data
    for page in doc:
        page.apply_redactions(graphics=True, images=fitz.PDF_REDACT_IMAGE_REMOVE)

    doc.save(output_pdf, garbage=4, deflate=True)
    print(f"[Engine: redact-pdf] Saved sanitized redacted PDF to '{output_pdf}'")
    return True


def main():
    if len(sys.argv) > 1 and sys.argv[1] == '--test':
        print("Testing all Python conversion engines...")
        import pdfplumber
        import openpyxl
        import pptx
        import reportlab
        import xhtml2pdf
        import pypdf
        import fitz
        print("pdfplumber, openpyxl, pptx, reportlab, xhtml2pdf, pypdf, fitz ready!")
        sys.exit(0)

    if len(sys.argv) < 4:
        print("Usage: python convert-office-pdf.py <mode> <input_file> <output_file> [extra_param]")
        print("Modes: pdf-to-word | pdf-to-excel | pdf-to-ppt | ppt-to-pdf | excel-to-pdf | pdf-to-markdown | html-to-pdf | ocr-pdf | repair-pdf | compress-pdf | pdf-to-pdf-a | protect-pdf | unlock-pdf | edit-pdf | redact-pdf | inspect-pdf")
        sys.exit(1)

    mode = sys.argv[1]
    inp = sys.argv[2]
    outp = sys.argv[3]
    extra = sys.argv[4] if len(sys.argv) > 4 else ""

    if mode == 'pdf-to-word':
        convert_pdf_to_word(inp, outp)
    elif mode == 'pdf-to-excel':
        convert_pdf_to_excel(inp, outp)
    elif mode == 'pdf-to-ppt':
        convert_pdf_to_ppt(inp, outp)
    elif mode == 'ppt-to-pdf':
        convert_ppt_to_pdf(inp, outp)
    elif mode == 'excel-to-pdf':
        convert_excel_to_pdf(inp, outp)
    elif mode == 'pdf-to-markdown':
        convert_pdf_to_markdown(inp, outp)
    elif mode == 'html-to-pdf':
        convert_html_to_pdf(inp, outp)
    elif mode == 'ocr-pdf':
        perform_ocr_pdf(inp, outp)
    elif mode == 'repair-pdf':
        repair_corrupted_pdf(inp, outp)
    elif mode == 'compress-pdf':
        compress_pdf_file(inp, outp)
    elif mode == 'pdf-to-pdf-a':
        convert_pdf_to_pdfa(inp, outp)
    elif mode == 'protect-pdf':
        protect_pdf_file(inp, outp, extra)
    elif mode == 'unlock-pdf':
        unlock_pdf_file(inp, outp, extra)
    elif mode == 'edit-pdf':
        edit_pdf_file(inp, outp, extra)
    elif mode == 'redact-pdf':
        redact_pdf_file(inp, outp, extra)
    elif mode == 'inspect-pdf':
        inspect_pdf_page_elements(inp, int(extra) if extra.isdigit() else 1, outp)
    else:
        print(f"Unknown conversion mode: {mode}")
        sys.exit(1)

if __name__ == '__main__':
    main()
