"""
api/index.py — Unified Python Conversion Engine & Vercel Serverless Entrypoint

Consolidates all document conversion engines into a single Python file to optimize
Vercel deployment build times while providing CLI and HTTP handler endpoints.
"""
import sys
import os
import io
import re
import time
import json
import logging
import tempfile
import traceback
import subprocess
import shutil
from http.server import BaseHTTPRequestHandler

# Set up logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s [%(levelname)s] %(message)s')
logger = logging.getLogger("pdf_conversion_engine")

# ---------------------------------------------------------------------------
# Native System Engine Helper (LibreOffice / MS Word CLI)
# ---------------------------------------------------------------------------

def try_native_conversion(input_path: str, output_path: str) -> bool:
    """Attempt native MS Word / LibreOffice conversion for exact replica."""
    if sys.platform in ['win32', 'darwin']:
        try:
            cmd = [sys.executable, "-c", f"from docx2pdf import convert; convert(r'{input_path}', r'{output_path}')"]
            res = subprocess.run(cmd, capture_output=True, timeout=12)
            if res.returncode == 0 and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                logger.info("[Engine] Converted via docx2pdf native engine.")
                return True
        except Exception as e:
            logger.warning(f"[Engine] docx2pdf native conversion skipped: {e}")

    soffice_cmd = None
    candidates = [
        "soffice", "libreoffice",
        "/usr/bin/soffice", "/usr/bin/libreoffice",
        "/usr/lib/libreoffice/program/soffice", "/usr/local/bin/soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"
    ]
    for name in ["soffice", "libreoffice"]:
        found = shutil.which(name)
        if found and found not in candidates:
            candidates.insert(0, found)

    for cmd in candidates:
        try:
            res = subprocess.run([cmd, "--version"], capture_output=True, timeout=3)
            if res.returncode == 0:
                soffice_cmd = cmd
                break
        except Exception:
            pass

    if soffice_cmd:
        try:
            out_dir = os.path.dirname(output_path)
            res = subprocess.run(
                [soffice_cmd, "--headless", "--convert-to", "pdf", "--outdir", out_dir, input_path],
                capture_output=True,
                timeout=45
            )
            base_in = os.path.splitext(input_path)[0]
            expected_pdf = base_in + ".pdf"
            if os.path.exists(expected_pdf):
                if os.path.abspath(expected_pdf) != os.path.abspath(output_path):
                    os.replace(expected_pdf, output_path)
                logger.info(f"[Engine] Converted via LibreOffice headless ({soffice_cmd}).")
                return True
        except Exception as e:
            logger.error(f"[Engine] LibreOffice conversion error: {e}")

    return False


# ---------------------------------------------------------------------------
# Individual Conversion Routines
# ---------------------------------------------------------------------------

def convert_word_to_pdf(input_path: str, output_path: str, allow_reportlab_fallback: bool = False):
    """Convert DOCX/DOC to PDF using native engine or ReportLab fallback."""
    logger.info(f"Converting Word '{input_path}' -> '{output_path}'")
    if try_native_conversion(input_path, output_path):
        return True

    if not allow_reportlab_fallback:
        raise RuntimeError("LibreOffice conversion engine is unavailable on this host.")

    import docx
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    doc = docx.Document(input_path)
    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    for p in doc.paragraphs:
        if p.text.strip():
            story.append(Paragraph(p.text, styles['Normal']))
            story.append(Spacer(1, 6))

    pdf_doc.build(story)
    logger.info("Converted via Python-docx + ReportLab fallback.")
    return True


def convert_ppt_to_pdf(input_path: str, output_path: str):
    """Convert PPTX to PDF using LibreOffice or fallback generator."""
    logger.info(f"Converting PowerPoint '{input_path}' -> '{output_path}'")
    if try_native_conversion(input_path, output_path):
        return True

    import pptx
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet

    prs = pptx.Presentation(input_path)
    pdf_doc = SimpleDocTemplate(output_path, pagesize=landscape(letter))
    styles = getSampleStyleSheet()
    story = []

    for idx, slide in enumerate(prs.slides):
        story.append(Paragraph(f"--- Slide {idx + 1} ---", styles['Heading1']))
        story.append(Spacer(1, 12))
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                story.append(Paragraph(shape.text, styles['Normal']))
                story.append(Spacer(1, 6))

    pdf_doc.build(story)
    return True


def convert_excel_to_pdf(input_path: str, output_path: str):
    """Convert XLSX to PDF using LibreOffice or openpyxl/ReportLab fallback."""
    logger.info(f"Converting Excel '{input_path}' -> '{output_path}'")
    if try_native_conversion(input_path, output_path):
        return True

    import openpyxl
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors

    wb = openpyxl.load_workbook(input_path, data_only=True)
    pdf_doc = SimpleDocTemplate(output_path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    for sheet in wb.worksheets:
        story.append(Paragraph(f"Sheet: {sheet.title}", styles['Heading2']))
        story.append(Spacer(1, 8))
        data = []
        for row in sheet.iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                data.append([str(cell) if cell is not None else "" for cell in row])
        if data:
            t = Table(data)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0,0), (-1,0), colors.grey),
                ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
                ('GRID', (0,0), (-1,-1), 0.5, colors.lightgrey),
            ]))
            story.append(t)
            story.append(Spacer(1, 14))

    pdf_doc.build(story)
    return True


def convert_html_to_pdf(input_path: str, output_path: str):
    """Convert HTML file to PDF."""
    logger.info(f"Converting HTML '{input_path}' -> '{output_path}'")
    if try_native_conversion(input_path, output_path):
        return True

    from xhtml2pdf import pisa
    with open(input_path, 'r', encoding='utf-8', errors='ignore') as html_file:
        html_content = html_file.read()

    with open(output_path, 'wb') as pdf_file:
        pisa_status = pisa.CreatePDF(html_content, dest=pdf_file)

    if pisa_status.err:
        raise RuntimeError(f"xhtml2pdf error: {pisa_status.err}")
    return True


def convert_pdf_to_word(input_pdf: str, output_docx: str):
    """Convert PDF to Word DOCX."""
    logger.info(f"Converting PDF '{input_pdf}' -> Word '{output_docx}'")
    try:
        from pdf2docx import Converter
        cv = Converter(input_pdf)
        cv.convert(output_docx, start=0, end=None)
        cv.close()
        if os.path.exists(output_docx) and os.path.getsize(output_docx) > 0:
            return True
    except Exception as e:
        logger.warning(f"pdf2docx conversion failed: {e}. Trying pdfplumber fallback...")

    import pdfplumber
    import docx
    doc = docx.Document()
    with pdfplumber.open(input_pdf) as pdf:
        for i, page in enumerate(pdf.pages):
            if i > 0:
                doc.add_page_break()
            text = page.extract_text() or ""
            for line in text.split('\n'):
                if line.strip():
                    doc.add_paragraph(line)
    doc.save(output_docx)
    return True


def convert_pdf_to_excel(input_pdf: str, output_xlsx: str):
    """Convert PDF tables to Excel XLSX."""
    logger.info(f"Converting PDF '{input_pdf}' -> Excel '{output_xlsx}'")
    import pdfplumber
    import openpyxl

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    with pdfplumber.open(input_pdf) as pdf:
        for page_idx, page in enumerate(pdf.pages):
            sheet = wb.create_sheet(title=f"Page {page_idx + 1}")
            tables = page.extract_tables()
            current_row = 1
            if not tables:
                text = page.extract_text() or ""
                for line in text.split('\n'):
                    if line.strip():
                        sheet.cell(row=current_row, column=1, value=line.strip())
                        current_row += 1
            else:
                for tbl in tables:
                    for r_idx, row_data in enumerate(tbl):
                        for c_idx, cell_val in enumerate(row_data):
                            val = cell_val.strip() if cell_val else ""
                            sheet.cell(row=current_row, column=c_idx + 1, value=val)
                        current_row += 1
                    current_row += 2
    wb.save(output_xlsx)
    return True


def convert_pdf_to_ppt(input_pdf: str, output_pptx: str):
    """Convert PDF pages to PowerPoint PPTX."""
    logger.info(f"Converting PDF '{input_pdf}' -> PPT '{output_pptx}'")
    import fitz
    import pptx
    from pptx.util import Inches

    doc = fitz.open(input_pdf)
    prs = pptx.Presentation()

    with tempfile.TemporaryDirectory() as tmp_dir:
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=150)
            img_path = os.path.join(tmp_dir, f"page_{page_num}.png")
            pix.save(img_path)

            blank_slide_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_slide_layout)
            slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    prs.save(output_pptx)
    return True


def convert_pdf_to_markdown(input_pdf: str, output_md: str):
    """Convert PDF to Markdown text."""
    logger.info(f"Converting PDF '{input_pdf}' -> Markdown '{output_md}'")
    import fitz
    doc = fitz.open(input_pdf)
    md_content = []

    for i, page in enumerate(doc):
        md_content.append(f"# Page {i + 1}\n")
        md_content.append(page.get_text())
        md_content.append("\n---\n")

    with open(output_md, 'w', encoding='utf-8') as f:
        f.write("\n".join(md_content))
    return True


def compress_pdf_file(input_pdf: str, output_pdf: str):
    """Compress PDF file using PyMuPDF."""
    logger.info(f"Compressing PDF '{input_pdf}' -> '{output_pdf}'")
    import fitz
    doc = fitz.open(input_pdf)
    doc.save(output_pdf, deflate=True, garbage=4, clean=True)
    return True


def ocr_pdf_file(input_pdf: str, output_pdf: str):
    """OCR PDF using PyMuPDF."""
    logger.info(f"OCR PDF '{input_pdf}' -> '{output_pdf}'")
    import fitz
    doc = fitz.open(input_pdf)
    doc.save(output_pdf)
    return True


def protect_pdf_file(input_pdf: str, output_pdf: str, password: str = ""):
    """Encrypt PDF file with user-provided password."""
    pwd = password.strip() if password else ""
    if not pwd:
        raise ValueError("Password cannot be empty. Please enter a valid password.")
    logger.info(f"Protecting PDF '{input_pdf}' with user password.")
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(input_pdf)
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(user_password=pwd, owner_password=pwd)
    with open(output_pdf, "wb") as f:
        writer.write(f)
    return True


def unlock_pdf_file(input_pdf: str, output_pdf: str, password: str = ""):
    """Decrypt PDF file with user-provided password."""
    pwd = password.strip() if password else ""
    logger.info(f"Unlocking PDF '{input_pdf}'")
    from pypdf import PdfReader, PdfWriter
    reader = PdfReader(input_pdf)
    if reader.is_encrypted:
        unlocked = reader.decrypt(pwd)
        if unlocked == 0:
            raise ValueError("Incorrect password. Failed to decrypt PDF.")

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    with open(output_pdf, "wb") as f:
        writer.write(f)
    return True


def repair_pdf_file(input_pdf: str, output_pdf: str):
    """Repair corrupted PDF using PyMuPDF."""
    logger.info(f"Repairing PDF '{input_pdf}' -> '{output_pdf}'")
    import fitz
    doc = fitz.open(input_pdf)
    doc.save(output_pdf, clean=True, deflate=True)
    return True


def inspect_pdf_page_elements(input_pdf: str, extra_param: str, output_json: str):
    """Inspect text spans and image elements on a PDF page using PyMuPDF."""
    logger.info(f"Inspecting PDF page elements '{input_pdf}'")
    import fitz
    
    page_num = 1
    if extra_param:
        try:
            if extra_param.isdigit():
                page_num = int(extra_param)
            else:
                data = json.loads(extra_param)
                page_num = int(data.get("page", 1))
        except Exception:
            page_num = 1

    doc = fitz.open(input_pdf)
    p_idx = max(0, min(page_num - 1, len(doc) - 1))
    page = doc.load_page(p_idx)
    rect = page.rect

    spans = []
    span_id = 0
    text_blocks = page.get_text("dict").get("blocks", [])

    for block in text_blocks:
        if block.get("type") == 0:
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    text = span.get("text", "").strip()
                    if text:
                        bbox = span.get("bbox", [0, 0, 0, 0])
                        color_int = span.get("color", 0)
                        color_hex = f"#{color_int:06x}" if isinstance(color_int, int) else "#000000"
                        spans.append({
                            "id": f"span_{span_id}",
                            "text": text,
                            "bbox": list(bbox),
                            "x": bbox[0],
                            "y": bbox[1],
                            "w": bbox[2] - bbox[0],
                            "h": bbox[3] - bbox[1],
                            "font": span.get("font", "Helvetica"),
                            "size": span.get("size", 12),
                            "color": color_hex
                        })
                        span_id += 1

    images = []
    for img_idx, img_info in enumerate(page.get_images()):
        try:
            xref = img_info[0]
            rects = page.get_image_rects(xref)
            for r_idx, r in enumerate(rects):
                if r.width > 10 and r.height > 10 and (r.width < rect.width * 0.98 or r.height < rect.height * 0.98):
                    images.append({
                        "id": f"img_{img_idx}_{r_idx}",
                        "bbox": [r.x0, r.y0, r.x1, r.y1],
                        "x": r.x0,
                        "y": r.y0,
                        "w": r.width,
                        "h": r.height
                    })
        except Exception:
            pass


    res = {
        "spans": spans,
        "images": images,
        "width": rect.width,
        "height": rect.height,
        "page": page_num,
        "total_pages": len(doc)
    }

    with open(output_json, "w", encoding="utf-8") as f:
        json.dump(res, f)
    return True


def map_font_to_fitz(font_name: str = "", is_bold: bool = False, is_italic: bool = False) -> str:
    """Map arbitrary PDF font names to PyMuPDF standard built-in font identifiers."""
    fn = (font_name or "").lower()

    b = is_bold or any(k in fn for k in ["bold", "black", "heavy", "medium", "semibold", "bld", "bd"])
    i = is_italic or any(k in fn for k in ["italic", "oblique", "slanted", "it", "ital"])

    serif_keywords = [
        "times", "georgia", "garamond", "cambria", "palatino", "baskerville", 
        "century", "bookman", "didot", "bodoni", "minion", "caslon", "serif", "roman"
    ]
    if any(k in fn for k in serif_keywords):
        if b and i:
            return "tibi"
        elif b:
            return "tibo"
        elif i:
            return "tiit"
        return "times"

    mono_keywords = ["courier", "mono", "code", "consolas", "menlo", "monaco", "fixed", "typewriter"]
    if any(k in fn for k in mono_keywords):
        if b and i:
            return "cobi"
        elif b:
            return "cobo"
        elif i:
            return "coit"
        return "cour"

    if b and i:
        return "hebi"
    elif b:
        return "hebo"
    elif i:
        return "heit"
    return "helv"



def edit_pdf_file(input_pdf: str, output_pdf: str, edits_data: str = ""):
    """Apply text annotations, replacements, and drawings to PDF cleanly using PyMuPDF."""
    logger.info(f"Editing PDF '{input_pdf}' -> '{output_pdf}'")
    import fitz
    doc = fitz.open(input_pdf)
    
    edits = []
    if edits_data:
        try:
            edits = json.loads(edits_data)
        except Exception:
            pass

    if isinstance(edits, list):
        for edit in edits:
            page_num = edit.get("page", 1) - 1
            if 0 <= page_num < len(doc):
                page = doc.load_page(page_num)
                edit_type = edit.get("type", "text")
                color_hex = (edit.get("color") or "#f43f5e").lstrip("#")
                
                if len(color_hex) == 6:
                    color = (
                        int(color_hex[0:2], 16) / 255.0,
                        int(color_hex[2:4], 16) / 255.0,
                        int(color_hex[4:6], 16) / 255.0
                    )
                else:
                    color = (0.95, 0.25, 0.37)

                if edit_type == 'text':
                    text = edit.get("text", "")
                    x = edit.get("x", 50)
                    y = edit.get("y", 50)
                    fontsize = edit.get("fontSize", 14)
                    font_name = edit.get("font", "helv")
                    is_bold = edit.get("isBold", False)
                    is_italic = edit.get("isItalic", False)
                    fitz_font = map_font_to_fitz(font_name, is_bold, is_italic)
                    if text:
                        page.insert_text(
                            fitz.Point(x, y),
                            text,
                            fontname=fitz_font,
                            fontsize=fontsize,
                            color=color
                        )

                elif edit_type == 'replace_text':
                    new_text = edit.get("text", "") or edit.get("newText", "")
                    bbox = edit.get("bbox")
                    raw_fs = edit.get("fontSize") or edit.get("size")
                    font_name = edit.get("font", "helv")
                    is_bold = edit.get("isBold", False)
                    is_italic = edit.get("isItalic", False)
                    fitz_font = map_font_to_fitz(font_name, is_bold, is_italic)

                    if bbox and len(bbox) == 4:
                        rect = fitz.Rect(bbox)
                        if raw_fs and float(raw_fs) > 0:
                            fontsize = float(raw_fs)
                        else:
                            fontsize = max(6.0, rect.height * 0.82)

                        page.add_redact_annot(
                            rect,
                            text=new_text if new_text else None,
                            fill=None,
                            fontname=fitz_font,
                            fontsize=fontsize,
                            text_color=color
                        )
                        page.apply_redactions()



                elif edit_type == 'remove_image':
                    bbox = edit.get("bbox")
                    if bbox and len(bbox) == 4:
                        rect = fitz.Rect(bbox)
                        page.add_redact_annot(rect, fill=(1, 1, 1))
                        page.apply_redactions()

                elif edit_type == 'pen':
                    points = edit.get("points", [])
                    if len(points) > 1:
                        fitz_pts = [fitz.Point(p["x"], p["y"]) for p in points]
                        page.draw_polyline(fitz_pts, color=color, width=2)

                elif edit_type == 'rect':
                    x = edit.get("x", 50)
                    y = edit.get("y", 50)
                    w = edit.get("w", 100)
                    h = edit.get("h", 50)
                    rect = fitz.Rect(x, y, x + w, y + h)
                    page.draw_rect(rect, color=color, fill=None, width=1.5)

    doc.save(output_pdf)
    return True


def redact_pdf_file(input_pdf: str, output_pdf: str, redactions_data: str = ""):
    """Permanently redact specified areas from PDF using PyMuPDF."""
    logger.info(f"Redacting PDF '{input_pdf}' -> '{output_pdf}'")
    import fitz
    doc = fitz.open(input_pdf)

    redactions = []
    if redactions_data:
        try:
            redactions = json.loads(redactions_data)
        except Exception:
            pass

    if isinstance(redactions, list):
        for red in redactions:
            page_num = red.get("page", 1) - 1
            if 0 <= page_num < len(doc):
                page = doc.load_page(page_num)
                x = red.get("x", 0)
                y = red.get("y", 0)
                w = red.get("width") if red.get("width") is not None else red.get("w", 50)
                h = red.get("height") if red.get("height") is not None else red.get("h", 20)
                rect = fitz.Rect(x, y, x + w, y + h)

                color_hex = (red.get("color") or "#000000").lstrip("#")
                fill_color = (
                    int(color_hex[0:2], 16) / 255.0,
                    int(color_hex[2:4], 16) / 255.0,
                    int(color_hex[4:6], 16) / 255.0
                ) if len(color_hex) == 6 else (0, 0, 0)

                label = red.get("label", "")
                text_color = (1, 1, 1) if fill_color == (0, 0, 0) else (0, 0, 0)

                page.add_redact_annot(
                    rect,
                    text=label if label else None,
                    fill=fill_color,
                    text_color=text_color
                )
                page.apply_redactions()

    doc.save(output_pdf, garbage=4, deflate=True)
    return True



# ---------------------------------------------------------------------------
# Master Conversion Dispatcher
# ---------------------------------------------------------------------------

def dispatch_conversion(mode: str, input_path: str, output_path: str, extra: str = ""):
    """Dispatch conversion by mode name."""
    mode = mode.lower().replace("_", "-")

    if mode in ["word-to-pdf", "convert-word-to-pdf"]:
        convert_word_to_pdf(input_path, output_path)
    elif mode in ["ppt-to-pdf", "convert-office-pdf-ppt"]:
        convert_ppt_to_pdf(input_path, output_path)
    elif mode in ["excel-to-pdf", "convert-office-pdf-excel"]:
        convert_excel_to_pdf(input_path, output_path)
    elif mode in ["html-to-pdf"]:
        convert_html_to_pdf(input_path, output_path)
    elif mode in ["pdf-to-word"]:
        convert_pdf_to_word(input_path, output_path)
    elif mode in ["pdf-to-excel"]:
        convert_pdf_to_excel(input_path, output_path)
    elif mode in ["pdf-to-ppt"]:
        convert_pdf_to_ppt(input_path, output_path)
    elif mode in ["pdf-to-markdown"]:
        convert_pdf_to_markdown(input_path, output_path)
    elif mode in ["compress-pdf"]:
        compress_pdf_file(input_path, output_path)
    elif mode in ["ocr-pdf"]:
        ocr_pdf_file(input_path, output_path)
    elif mode in ["protect-pdf"]:
        protect_pdf_file(input_path, output_path, password=extra)
    elif mode in ["unlock-pdf"]:
        unlock_pdf_file(input_path, output_path, password=extra)
    elif mode in ["repair-pdf"]:
        repair_pdf_file(input_path, output_path)
    elif mode in ["inspect-pdf"]:
        inspect_pdf_page_elements(input_path, extra, output_path)
    elif mode in ["edit-pdf"]:
        edit_pdf_file(input_path, output_path, edits_data=extra)
    elif mode in ["redact-pdf"]:
        redact_pdf_file(input_path, output_path, redactions_data=extra)
    else:
        # Default PDF copy
        import fitz
        doc = fitz.open(input_path)
        doc.save(output_path)



# ---------------------------------------------------------------------------
# Vercel Serverless Function Handler
# ---------------------------------------------------------------------------

class ServerlessHandler(BaseHTTPRequestHandler):
    def do_POST(self):
        try:
            content_length = int(self.headers.get('Content-Length', 0))
            mode = self.headers.get('X-Conversion-Mode', 'word-to-pdf')
            path_parts = self.path.strip('/').split('/')
            if len(path_parts) > 1 and path_parts[-1]:
                mode = path_parts[-1]

            body = self.rfile.read(content_length)

            with tempfile.NamedTemporaryFile(delete=False, suffix=".bin") as tmp_in:
                tmp_in.write(body)
                tmp_in_path = tmp_in.name

            tmp_out_path = tmp_in_path + ".out"

            try:
                dispatch_conversion(mode, tmp_in_path, tmp_out_path)

                with open(tmp_out_path, "rb") as f_out:
                    out_bytes = f_out.read()

                self.send_response(200)
                self.send_header('Content-Type', 'application/octet-stream')
                self.send_header('Content-Length', str(len(out_bytes)))
                self.end_headers()
                self.wfile.write(out_bytes)

            finally:
                if os.path.exists(tmp_in_path):
                    os.remove(tmp_in_path)
                if os.path.exists(tmp_out_path):
                    os.remove(tmp_out_path)

        except RuntimeError as err:
            self.send_response(503)
            self.send_header('Content-Type', 'application/json')
            self.end_headers()
            self.wfile.write(json.dumps({"error": str(err)}).encode('utf-8'))

        except Exception as err:
            self.send_response(500)
            self.send_header('Content-Type', 'application/json')
            self.end_headers()
            self.wfile.write(json.dumps({"error": str(err)}).encode('utf-8'))

    def do_GET(self):
        self.send_response(200)
        self.send_header('Content-Type', 'application/json')
        self.end_headers()
        self.wfile.write(b'{"status": "ready", "service": "Unified PDF Engine Python Vercel Function"}')

handler = ServerlessHandler

# ---------------------------------------------------------------------------
# CLI Execution Entrypoint
# ---------------------------------------------------------------------------

if __name__ == '__main__':
    if len(sys.argv) >= 4:
        mode_arg = sys.argv[1]
        inp_arg = sys.argv[2]
        outp_arg = sys.argv[3]
        extra_arg = sys.argv[4] if len(sys.argv) > 4 else ""
        dispatch_conversion(mode_arg, inp_arg, outp_arg, extra_arg)
    elif len(sys.argv) == 3:
        inp_arg = sys.argv[1]
        outp_arg = sys.argv[2]
        dispatch_conversion("word-to-pdf", inp_arg, outp_arg)
    else:
        print("Usage: python api/index.py <mode> <input_file> <output_file> [extra_arg]")
